"""PDF → PPT — 将 PDF 每页转为 PPT 幻灯片"""
import os
import threading
from typing import Optional

import fitz
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

from src.core.converter_base import BaseConverter
from src.core.converter_registry import register
from src.core.models import ConversionType

# PPT 标准尺寸（英寸）
SLIDE_WIDTH_INCHES = 13.333
SLIDE_HEIGHT_INCHES = 7.5


@register(ConversionType.PDF_TO_PPT)
class PdfToPptConverter(BaseConverter):
    """将 PDF 转换为 PPT 演示文稿"""

    def convert(
        self,
        input_path: str,
        output_dir: str,
        options: Optional[dict] = None,
        cancel_event: Optional[threading.Event] = None,
    ) -> str:
        options = options or {}
        dpi = options.get("dpi", 200)

        doc = fitz.open(input_path)
        prs = Presentation()

        # 设置幻灯片尺寸为宽屏 16:9
        prs.slide_width = Inches(SLIDE_WIDTH_INCHES)
        prs.slide_height = Inches(SLIDE_HEIGHT_INCHES)

        for page_num, page in enumerate(doc):
            self._check_cancelled(cancel_event)

            # 使用空白布局
            blank_layout = prs.slide_layouts[6]  # blank
            slide = prs.slides.add_slide(blank_layout)

            # 1. 将 PDF 页面渲染为图片作为幻灯片背景
            pix = page.get_pixmap(dpi=150)  # 适中 DPI
            img_bytes = pix.tobytes("png")

            # 插入页面截图
            from io import BytesIO
            img_stream = BytesIO(img_bytes)

            # 计算图片在幻灯片中的适配
            slide_w = prs.slide_width
            slide_h = prs.slide_height
            img_w = Emu(int(pix.width * 914400 / 150))  # 转换为 EMU
            img_h = Emu(int(pix.height * 914400 / 150))

            # 按比例缩放以适应幻灯片
            scale_w = slide_w / img_w
            scale_h = slide_h / img_h
            scale = min(scale_w, scale_h)

            final_w = Emu(int(img_w * scale))
            final_h = Emu(int(img_h * scale))
            left = Emu(int((slide_w - final_w) / 2))
            top = Emu(int((slide_h - final_h) / 2))

            slide.shapes.add_picture(img_stream, left, top, final_w, final_h)

            # 2. 提取文本块并添加为文本框
            blocks = page.get_text("dict")["blocks"]
            for block in blocks:
                if block["type"] != 0:  # 跳过图片块
                    continue
                for line in block.get("lines", []):
                    text_parts = []
                    for span in line.get("spans", []):
                        text = span.get("text", "").strip()
                        if text:
                            text_parts.append(text)
                    if not text_parts:
                        continue

                    full_text = " ".join(text_parts)
                    bbox = line["bbox"]  # PDF 坐标 (x0, y0, x1, y1)

                    # 将 PDF 坐标映射到幻灯片坐标
                    pdf_w = page.rect.width
                    pdf_h = page.rect.height

                    left_emu = Emu(int(bbox[0] * slide_w / pdf_w))
                    top_emu = Emu(int(bbox[1] * slide_h / pdf_h))
                    width_emu = Emu(int((bbox[2] - bbox[0]) * slide_w / pdf_w * 1.2))
                    height_emu = Emu(int((bbox[3] - bbox[1]) * slide_h / pdf_h * 1.2))

                    # 获取字体信息
                    first_span = line["spans"][0] if line.get("spans") else {}
                    font_size_pt = first_span.get("size", 11)

                    txBox = slide.shapes.add_textbox(left_emu, top_emu, width_emu, height_emu)
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.text = full_text
                    p.font.size = Pt(min(font_size_pt, 24))

        doc.close()

        output_path = self._make_output_path(input_path, output_dir, ".pptx")
        prs.save(output_path)
        return output_path
